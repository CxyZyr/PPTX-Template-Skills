"""parse_template.py — turn a .pptx template into a semantic template_spec.json
(+ per-slide PNG renders for visual grounding).

Usage:
  python parse_template.py --template "PPT模板/大气人工智能科技感PPT模板.pptx" \
                           --out workspace/specs/ai_tech [--no-render] [--limit N] [--dpi 110]

Output:
  <out>/spec.json        the machine-readable contract consumed by pptx-generation
  <out>/renders/*.png    one image per slide (skipped with --no-render)
  <out>/summary.txt      human-readable role overview for quick review
"""
from __future__ import annotations

import argparse
import collections
import json
import pathlib
import shutil
import subprocess
import sys

SCRIPT_DIR = pathlib.Path(__file__).resolve().parent
SKILL_ROOT = SCRIPT_DIR.parents[0]
REPO_ROOT = SCRIPT_DIR.parents[2]
sys.path.insert(0, str(REPO_ROOT))
sys.path.insert(0, str(SKILL_ROOT))
sys.path.insert(0, str(SCRIPT_DIR))

from pptx import Presentation  # noqa: E402

from lib.pptx_toolkit import render as render_mod  # noqa: E402
from shape_classifier import classify_slide  # noqa: E402
from slide_classifier import analyze_slide  # noqa: E402

SCHEMA_VERSION = "1.0"


def serialize_node(node: dict) -> dict:
    g = node["geometry"]
    out = {
        "path": node["path"],
        "shape_type": node["shape_type"],
        "role": node.get("role", "unknown"),
        "role_confidence": round(node.get("role_confidence", 0.0), 2),
        "signals": node.get("signals", []),
        "geometry": {
            # EMU (generation needs these to resize/position)
            "left": g["left"], "top": g["top"], "width": g["width"], "height": g["height"],
            # readable
            "left_in": g["left_in"], "top_in": g["top_in"],
            "width_in": g["width_in"], "height_in": g["height_in"],
            "area_ratio": g["area_ratio"], "aspect": g["aspect"], "on_canvas": g["on_canvas"],
        },
    }
    if node.get("placeholder"):
        out["placeholder"] = node["placeholder"]
    if node.get("text"):
        t = node["text"]
        out["text"] = {
            "sample": t["sample"], "char_count": t["char_count"],
            "max_font_pt": t["max_font_pt"], "min_font_pt": t["min_font_pt"],
            "any_bold": t["any_bold"], "font_name": t.get("font_name"),
            "bold": t.get("bold"), "italic": t.get("italic"),
            "align": t["align"], "color": t["color"],
            "paragraphs": t.get("paragraphs", []),
            "vertical": t["vertical"],
        }
    if node["children"]:
        out["children"] = [serialize_node(c) for c in node["children"]]
    return out


def extract_theme(prs) -> dict:
    """Best-effort: most common run font + the palette of explicit run colors."""
    fonts = collections.Counter()
    colors = collections.Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts[run.font.name] += 1
                    try:
                        if run.font.color is not None and run.font.color.type is not None:
                            colors[str(run.font.color.rgb)] += 1
                    except Exception:
                        pass
    common_fonts = [f for f, _ in fonts.most_common(4)]
    return {
        "common_fonts": common_fonts,
        "font_availability": font_availability(common_fonts),
        "common_colors": [c for c, _ in colors.most_common(6)],
    }


def font_availability(fonts: list[str]) -> dict:
    """Best-effort installed-font check for render diagnostics."""
    if not fonts:
        return {}
    if shutil.which("fc-match") is None:
        return {font: {"available": None, "matched": None} for font in fonts}
    availability = {}
    for font in fonts:
        try:
            proc = subprocess.run(
                ["fc-match", "-f", "%{family}", font],
                check=False,
                text=True,
                capture_output=True,
                timeout=3,
            )
            matched = (proc.stdout or "").strip()
            matched_names = [part.strip().casefold() for part in matched.split(",") if part.strip()]
            availability[font] = {
                "available": font.casefold() in matched_names,
                "matched": matched or None,
            }
        except Exception:
            availability[font] = {"available": None, "matched": None}
    return availability


def build_spec(template_path: pathlib.Path, *, render: bool, out_dir: pathlib.Path,
               limit: int | None, dpi: int) -> dict:
    prs = Presentation(str(template_path))
    slide_w, slide_h = prs.slide_width, prs.slide_height
    total = len(prs.slides)

    renders: list = []
    if render:
        try:
            renders = render_mod.render_deck(template_path, out_dir / "renders", dpi=dpi)
        except Exception as exc:  # rendering is best-effort
            print(f"[warn] rendering failed ({exc}); continuing without renders", file=sys.stderr)
            renders = []

    slides_out = []
    for index, slide in enumerate(prs.slides):
        if limit is not None and index >= limit:
            break
        nodes = classify_slide(slide.shapes, slide_w, slide_h)
        analysis = analyze_slide(nodes, index, total)
        render_rel = None
        if index < len(renders):
            render_rel = str(pathlib.Path("renders") / renders[index].name)
        slides_out.append({
            "index": index,
            "role": analysis["role"],
            "role_confidence": round(analysis["role_confidence"], 2),
            "role_signals": analysis["role_signals"],
            "render": render_rel,
            "fill_plan": analysis["fill_plan"],
            "groups": analysis["groups"],
            "shapes": [serialize_node(n) for n in nodes],
        })

    return {
        "schema_version": SCHEMA_VERSION,
        "template": {
            "path": str(template_path),
            "name": template_path.stem,
            "slide_width_emu": slide_w,
            "slide_height_emu": slide_h,
            "slide_count": total,
            "theme": extract_theme(prs),
        },
        "slides": slides_out,
    }


def write_summary(spec: dict, path: pathlib.Path) -> None:
    lines = [f"# {spec['template']['name']}  ({spec['template']['slide_count']} slides)",
             f"theme fonts: {spec['template']['theme']['common_fonts']}", ""]
    font_status = spec["template"].get("theme", {}).get("font_availability") or {}
    missing = [
        f"{font} -> {info.get('matched')}"
        for font, info in font_status.items()
        if info.get("available") is False
    ]
    if missing:
        lines.append("missing fonts: " + "; ".join(missing))
        lines.append("")
    for s in spec["slides"]:
        fp = s["fill_plan"]
        bits = []
        if fp["title"]:
            bits.append(f"title='{fp['title']['sample']}'")
        if fp["cards"]:
            bits.append(f"cards={len(fp['cards'])}")
        if fp["body"]:
            bits.append(f"body={len(fp['body'])}")
        if fp.get("labels"):
            bits.append(f"labels={len(fp['labels'])}")
        if fp["images"]:
            bits.append(f"images={len(fp['images'])}")
        if fp.get("charts"):
            bits.append(f"charts={len(fp['charts'])}")
        if fp.get("tables"):
            bits.append(f"tables={len(fp['tables'])}")
        if fp["logo"]:
            bits.append("logo")
        if fp["section_number"]:
            bits.append(f"sec='{fp['section_number']['sample']}'")
        lines.append(f"slide {s['index']:02d}  [{s['role']:15s} {s['role_confidence']:.2f}]  "
                     + "  ".join(bits))
    path.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    ap = argparse.ArgumentParser(description="Parse a PPTX template into a semantic spec.")
    ap.add_argument("--template", required=True)
    ap.add_argument("--out", required=True)
    ap.add_argument("--no-render", action="store_true", help="skip slide rendering")
    ap.add_argument("--limit", type=int, default=None, help="only parse first N slides")
    ap.add_argument("--dpi", type=int, default=110)
    args = ap.parse_args()

    template_path = pathlib.Path(args.template)
    out_dir = pathlib.Path(args.out)
    out_dir.mkdir(parents=True, exist_ok=True)

    spec = build_spec(
        template_path, render=not args.no_render, out_dir=out_dir,
        limit=args.limit, dpi=args.dpi,
    )
    (out_dir / "spec.json").write_text(json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8")
    write_summary(spec, out_dir / "summary.txt")
    print(f"spec  -> {out_dir / 'spec.json'}")
    print(f"summary -> {out_dir / 'summary.txt'}")
    print((out_dir / "summary.txt").read_text(encoding="utf-8"))


if __name__ == "__main__":
    main()
