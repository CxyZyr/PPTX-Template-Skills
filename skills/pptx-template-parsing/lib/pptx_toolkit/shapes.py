"""Shape manipulation: delete, hide, and replace shapes with pictures.

Extracted from scripts/generate_strategy_ppt_ai_tech.py and
scripts/generate_gcl_strategy_ppt_tech_release.py.
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def delete_shape(shape) -> None:
    """Remove a shape from its parent (slide or group)."""
    shape._element.getparent().remove(shape._element)  # noqa: SLF001


def hide_shape_visual(shape) -> None:
    """Make a shape visually disappear (clear text, transparent fill + line),
    recursing into groups — without deleting it (keeps the layout skeleton)."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            hide_shape_visual(child)
        return

    if getattr(shape, "has_text_frame", False):
        shape.text_frame.clear()

    if hasattr(shape, "fill"):
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _WHITE
            shape.fill.transparency = 1.0
        except Exception:
            pass

    if hasattr(shape, "line"):
        try:
            shape.line.transparency = 1.0
        except Exception:
            pass


def add_centered_picture(slide, image_path, left: int, top: int, width: int, height: int,
                         *, padding_ratio: float = 0.18):
    """Place an image centered inside the (left, top, width, height) EMU box,
    preserving the image aspect ratio and leaving `padding_ratio` margin."""
    image_path = str(image_path)
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    avail_w = max(1, int(width * (1 - padding_ratio)))
    avail_h = max(1, int(height * (1 - padding_ratio)))
    scale = min(avail_w / img_w, avail_h / img_h)
    pic_w = max(1, int(img_w * scale))
    pic_h = max(1, int(img_h * scale))
    pic_left = left + (width - pic_w) // 2
    pic_top = top + (height - pic_h) // 2
    return slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_w, height=pic_h)


def replace_shape_with_picture(slide, shape, image_path, *, padding_ratio: float = 0.12,
                               centered: bool = True):
    """Replace `shape` with a picture occupying its box, then delete the shape.

    centered=True keeps the image aspect ratio centered in the box (good for
    logos/icons). centered=False stretches to fill the box exactly (good when
    the image is already cropped to the box aspect ratio)."""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    if centered:
        pic = add_centered_picture(slide, image_path, left, top, width, height,
                                   padding_ratio=padding_ratio)
    else:
        pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    delete_shape(shape)
    return pic
