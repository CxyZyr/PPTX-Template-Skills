"""Content application: fill one slide from a content_plan `page` against a
parsed `spec.json` `fill_plan`.

This is the generic applier that makes the parse -> generate pipeline real.
The parsing skill (pptx-template-parsing) emits, per slide, a `fill_plan` that
maps semantic slots (title / subtitle / section_number / logo / footer /
labels[] / body[] / images[] / cards[]) to index-paths. The generation skill
(ppt-template-adaptation) authors a `page` dict that supplies content per slot
*by name*. `apply_page` resolves the paths and fills, reusing the proven lib
helpers (autofit/text/assets/shapes). No index-paths are hand-written upstream.

Design rules:
  - Resolve every slot path to a shape object in ONE read pass before mutating
    (deleting a sibling keeps held refs valid; re-resolving by index would not).
  - Never raise on a single bad/missing slot — collect warnings, keep going.
  - Unspecified style inherits the template (color/font None -> template run).
  - Card `residual_text` paths are parser-owned template text inside a card
    group that was not selected as title/body/number; clear them after filling
    the card so stale sample labels do not survive.
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .assets import crop_to_fill, get_icon
from .autofit import candidate_breaks, fit_body, fit_title
from .fonts import measure_line_height_pt, measure_line_width_pt
from .geometry import emu_to_pt, inch_to_emu, pt_to_emu
from .shapes import add_centered_picture, delete_shape, hide_shape_visual, replace_shape_with_picture
from .text import replace_first_run_text, set_text
from .walk import absolute_box_for_path, get_shape_by_path

# Stock placeholder text that must never survive into a delivered deck.
STOCK_TEXT_RE = re.compile(
    r"点击|单击|请在|在此|输入|添加|标题文字|文字内容|内容描述|详细内容|示例|"
    r"Lorem|YOUR\s*LOGO|您的\s*logo|公司\s*logo|\bLOGO\b",
    re.I,
)
LEADING_NUMBER_RE = re.compile(r"^\s*([0-9]{1,3}|[A-Z]\d{1,3})[\s.、:：\-_/]+(.+?)\s*$")
NUMBER_ONLY_RE = re.compile(r"^\s*([0-9]{1,3}|[A-Z]\d{1,3})\s*$")
BODYISH_TEXT_RE = re.compile(
    r"(点击|单击|请|在此).{0,6}(输入|添加).{0,4}(文字|文本|内容|正文|介绍)"
    r"|文字内容|详细内容|内容描述|输入内容|Please\s+enter\s+details",
    re.I,
)


# ---------------------------------------------------------------------------
# small readers / converters
# ---------------------------------------------------------------------------
def _to_rgb(value):
    if value is None:
        return None
    if isinstance(value, RGBColor):
        return value
    if isinstance(value, (list, tuple)) and len(value) == 3:
        return RGBColor(*value)
    if isinstance(value, dict):
        rgb = value.get("rgb")
        if rgb:
            return RGBColor.from_string(str(rgb).lstrip("#"))
        theme = str(value.get("theme_color") or "").upper()
        if theme in {"BACKGROUND_1", "LIGHT_1"}:
            return RGBColor(255, 255, 255)
        if theme in {"TEXT_1", "DARK_1"}:
            return RGBColor(0, 0, 0)
        return None
    return RGBColor.from_string(str(value).lstrip("#"))


def _rgb_tuple(value) -> tuple[int, int, int] | None:
    rgb = _to_rgb(value)
    if rgb is None:
        return None
    return int(rgb[0]), int(rgb[1]), int(rgb[2])


def _as_list(value) -> list:
    if value is None:
        return []
    return value if isinstance(value, list) else [value]


def _has_content(value) -> bool:
    return any(str(item).strip() for item in _as_list(value))


def _current_pt(shape, default: float) -> float:
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    return run.font.size.pt
    except Exception:
        pass
    return default


def _current_align(shape):
    try:
        for para in shape.text_frame.paragraphs:
            if para.alignment is not None:
                return para.alignment
    except Exception:
        pass
    return None


def _current_bold(shape) -> bool:
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.bold is not None:
                    return bool(run.font.bold)
    except Exception:
        pass
    return False


def _set_word_wrap(shape, value: bool) -> None:
    try:
        shape.text_frame.word_wrap = value
    except Exception:
        pass


def _align_from_value(value):
    if value is None:
        return None
    if isinstance(value, PP_ALIGN):
        return value
    normalized = str(value).strip().upper()
    return {
        "LEFT": PP_ALIGN.LEFT,
        "CENTER": PP_ALIGN.CENTER,
        "CENTRE": PP_ALIGN.CENTER,
        "RIGHT": PP_ALIGN.RIGHT,
        "JUSTIFY": PP_ALIGN.JUSTIFY,
        "DISTRIBUTE": PP_ALIGN.DISTRIBUTE,
    }.get(normalized)


def _clear_if_stock(shape) -> bool:
    """Blank a text slot iff it still holds template placeholder text.

    Keeps intentional decorative text, removes leftover `点击输入标题` etc.
    """
    if shape is None or not getattr(shape, "has_text_frame", False):
        return False
    text = shape.text_frame.text.strip()
    if text and STOCK_TEXT_RE.search(text):
        replace_first_run_text(shape, "")
        return True
    return False


def _norm_text(value: str | None) -> str:
    return re.sub(r"\s+", "", value or "")


def _clear_text(shape) -> bool:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return False
    for para in shape.text_frame.paragraphs:
        if para.runs:
            for run in para.runs:
                run.text = ""
        else:
            para.text = ""
    return True


def _clear_if_unfilled(shape, slot: dict | None = None) -> bool:
    """Clear an unfilled slot when the parser marked its sample as stale."""
    if shape is None or not getattr(shape, "has_text_frame", False):
        return False
    text = shape.text_frame.text.strip()
    if not text:
        return False
    if STOCK_TEXT_RE.search(text):
        return _clear_text(shape)
    slot = slot or {}
    if slot.get("default_action", "clear") != "clear":
        return False
    sample = slot.get("sample")
    if sample and _norm_text(text) == _norm_text(sample):
        return _clear_text(shape)
    return False


def _clear_stock_texts(shapes) -> int:
    """Recursively clear leftover template placeholder text missed by slots.

    The parser exposes known fill slots, but some templates keep extra prompt
    text inside decorative groups. A delivered deck should never retain stock
    authoring prompts, so this final pass removes only text that matches the
    generic placeholder regex.
    """
    cleared = 0
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text and STOCK_TEXT_RE.search(text):
                if _clear_text(shape):
                    cleared += 1
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            cleared += _clear_stock_texts(shape.shapes)
    return cleared


# ---------------------------------------------------------------------------
# resolve (read pass) — every slot path -> shape object, before any mutation
# ---------------------------------------------------------------------------
def resolve_fill_plan(slide, fill_plan: dict) -> dict:
    """Resolve all fill_plan slot paths to live shape objects (or None)."""
    def one(slot):
        return get_shape_by_path(slide, slot["path"]) if slot else None

    resolved = {
        "title": one(fill_plan.get("title")),
        "subtitle": one(fill_plan.get("subtitle")),
        "section_number": one(fill_plan.get("section_number")),
        "logo": one(fill_plan.get("logo")),
        "footer": one(fill_plan.get("footer")),
        "labels": [get_shape_by_path(slide, lab["path"]) for lab in fill_plan.get("labels", [])],
        "body": [get_shape_by_path(slide, b["path"]) for b in fill_plan.get("body", [])],
        "images": [
            {
                "shape": get_shape_by_path(slide, im["path"]),
                "box": absolute_box_for_path(slide, im["path"]),
            }
            for im in fill_plan.get("images", [])
        ],
        "cards": [],
    }
    for card in fill_plan.get("cards", []):
        resolved["cards"].append({
            "title": get_shape_by_path(slide, card["title"]) if card.get("title") else None,
            "title_box": absolute_box_for_path(slide, card["title"]) if card.get("title") else None,
            "body": get_shape_by_path(slide, card["body"]) if card.get("body") else None,
            "body_box": absolute_box_for_path(slide, card["body"]) if card.get("body") else None,
            "icon": get_shape_by_path(slide, card["icon"]) if card.get("icon") else None,
            "icon_box": absolute_box_for_path(slide, card["icon"]) if card.get("icon") else None,
            "icon_overlays": [
                get_shape_by_path(slide, path)
                for path in card.get("icon_overlays", [])
            ],
            "card_box": _geometry_box(card.get("card_geometry")),
            "icon_style": card.get("icon_style"),
            "number": get_shape_by_path(slide, card["number"]) if card.get("number") else None,
            "residual_text": [
                get_shape_by_path(slide, path)
                for path in card.get("residual_text", [])
            ],
            "title_body_mode": card.get("title_body_mode"),
            "title_style": card.get("title_style"),
            "body_style": card.get("body_style"),
        })
    return resolved


# ---------------------------------------------------------------------------
# individual fillers (operate on already-resolved shapes)
# ---------------------------------------------------------------------------
def _fit_title(shape, text, *, font, color, min_left=None, max_right=None,
               max_bottom=None, allow_break=False, align=None,
               prefer_break=False, respect_bounds=False, bold=None,
               min_size_ratio=0.6):
    base = round(_current_pt(shape, 28))
    set_align = align or _current_align(shape) or PP_ALIGN.LEFT
    set_bold = _current_bold(shape) if bold is None else bold
    fit_title(
        shape, text,
        base_size_pt=base, min_size_pt=max(12, int(base * min_size_ratio)),
        font_name=font, color=color, align=set_align, bold=set_bold,
        min_left=min_left, max_right=max_right, max_bottom=max_bottom,
        allow_break=allow_break, prefer_break=prefer_break,
        respect_bounds=respect_bounds,
    )
    _set_word_wrap(shape, False)


def _fit_label(shape, text, *, font, slide_width=None, color=None, visual_box=None,
               max_abs_right=None, max_abs_bottom=None):
    """Fit longer labels without forcing a line break.

    Short labels/numbers should keep exact template styling via run replacement.
    Longer footer/tagline labels need width/size adaptation or LibreOffice may
    wrap the final character even when the visual layout can tolerate a wider
    box.
    """
    base = round(_current_pt(shape, 16))
    align = _current_align(shape) or PP_ALIGN.CENTER
    bold = _current_bold(shape)
    min_left = inch_to_emu(0.35)
    max_right = (slide_width - inch_to_emu(0.35)) if slide_width else None
    if visual_box and slide_width:
        left, _top, abs_width, abs_height = visual_box
        right = left + abs_width
        if align in (PP_ALIGN.CENTER, PP_ALIGN.DISTRIBUTE):
            max_abs_width = max_right - min_left
        elif align == PP_ALIGN.RIGHT:
            max_abs_width = right - min_left
        else:
            max_abs_width = max_right - left
        if max_abs_right is not None:
            max_abs_width = min(max_abs_width, max_abs_right - left)
        max_abs_width = max(1, max_abs_width)
        current_abs_width = min(abs_width, max_abs_width)
        min_size = max(10, int(base * 0.55))
        best_size = min_size
        best_abs_width = current_abs_width
        found_fit = False
        for size_pt in range(base, min_size - 1, -1):
            line_width = measure_line_width_pt(text, size_pt, bold=bold)
            margin = max(24, size_pt * 1.1)
            needed_abs_width = pt_to_emu(line_width * 1.12 + margin)
            if needed_abs_width <= max_abs_width:
                best_size = size_pt
                best_abs_width = max(current_abs_width, needed_abs_width)
                found_fit = True
                break
        if not found_fit:
            best_abs_width = max(current_abs_width, max_abs_width)
        _resize_shape_to_visual_width(
            shape,
            visual_box,
            best_abs_width,
            align=align,
        )
        line_height = measure_line_height_pt(best_size, bold=bold)
        scale_y = (abs_height / shape.height) if shape.height else 1.0
        if scale_y:
            new_height = max(shape.height, int(pt_to_emu(line_height + 8) / scale_y))
            if max_abs_bottom is not None:
                local_limit = _local_limit_bottom(shape, visual_box, max_abs_bottom)
                new_height = min(new_height, max(1, local_limit - shape.top))
            shape.height = new_height
        set_text(
            shape,
            [text],
            font_name=font,
            size=Pt(best_size),
            bold=bold,
            color=color,
            align=align,
            line_spacing=1.0,
            space_after=0,
        )
        _set_word_wrap(shape, False)
        return
    fit_title(
        shape, text,
        base_size_pt=base,
        min_size_pt=max(10, int(base * 0.55)),
        font_name=font,
        color=color,
        align=align,
        bold=bold,
        min_left=min_left,
        max_right=max_right,
        allow_break=False,
    )
    _set_word_wrap(shape, False)


def _fit_bounded_title(shape, text, *, font, color, visual_box=None, bold=None):
    """Fit a repeated card-list title inside its visual box, allowing breaks."""
    set_bold = _current_bold(shape) if bold is None else bold
    if not visual_box:
        _fit_title(shape, text, font=font, color=color, allow_break=True, prefer_break=True, bold=set_bold)
        return
    base = round(_current_pt(shape, 28))
    min_size = max(12, int(base * 0.6))
    align = _current_align(shape) or PP_ALIGN.CENTER
    max_width_pt = max(20, emu_to_pt(visual_box[2]) - max(10, base * 0.22))
    max_height_pt = max(20, emu_to_pt(visual_box[3]))
    options = candidate_breaks(text) + [text]
    seen = set()
    best = None
    for candidate in options:
        if candidate in seen:
            continue
        seen.add(candidate)
        lines = [line.strip() for line in candidate.split("\n") if line.strip()]
        if not lines:
            continue
        for size_pt in range(base, min_size - 1, -2):
            line_width = max(measure_line_width_pt(line, size_pt, bold=set_bold) for line in lines)
            line_height = measure_line_height_pt(size_pt, bold=set_bold)
            needed_height = line_height * len(lines) * 1.02 + 4
            if line_width <= max_width_pt and needed_height <= max_height_pt * 1.18:
                best = (lines, size_pt)
                break
        if best:
            break
    if best is None:
        # Preserve the bounded layout even when fitting is tight: prefer a
        # natural break at the smallest allowed size over one long overlapping line.
        fallback = candidate_breaks(text)[0] if candidate_breaks(text) else text
        best = ([line.strip() for line in fallback.split("\n") if line.strip()], min_size)
    lines, size_pt = best
    set_text(
        shape,
        lines,
        font_name=font,
        size=Pt(size_pt),
        bold=set_bold,
        color=color,
        align=align,
        line_spacing=0.96,
        space_after=0,
    )
    _set_word_wrap(shape, False)


def _resize_shape_to_visual_width(shape, visual_box, target_abs_width: int, *, align) -> None:
    """Resize a shape by displayed width, preserving center/right anchors.

    Group children use a local coordinate system. `visual_box` is slide-level
    geometry; comparing it to local width gives the x-scale needed to resize
    without hand-writing group coordinates.
    """
    if shape is None or not visual_box:
        return
    current_abs_width = max(1, int(visual_box[2]))
    target_abs_width = max(current_abs_width, int(target_abs_width))
    if target_abs_width <= current_abs_width:
        return
    scale_x = current_abs_width / shape.width if shape.width else 1.0
    if not scale_x:
        return
    new_width = int(round(target_abs_width / scale_x))
    delta = new_width - shape.width
    if delta <= 0:
        return
    if align in (PP_ALIGN.CENTER, PP_ALIGN.DISTRIBUTE):
        shape.left = int(shape.left - delta / 2)
    elif align == PP_ALIGN.RIGHT:
        shape.left = int(shape.left - delta)
    shape.width = new_width


def _title_allows_break(shape, slide_width: int | None, *, right_constrained: bool = False) -> bool:
    """Only centered/display title boxes should invent a title line break.

    Left/top section titles should expand to the right first. This keeps
    `pdftotext` line wrapping from being mistaken for a layout failure.
    """
    if right_constrained:
        return True
    align = _current_align(shape)
    if align in (PP_ALIGN.CENTER, PP_ALIGN.DISTRIBUTE):
        return True
    if not slide_width:
        return False
    cx_ratio = (shape.left + shape.width / 2) / slide_width
    width_ratio = shape.width / slide_width
    return 0.36 <= cx_ratio <= 0.64 and width_ratio >= 0.34


def _right_limit_is_constrained(limits: dict, default_right: int | None) -> bool:
    if not default_right:
        return False
    max_abs_right = limits.get("max_abs_right")
    if max_abs_right is None:
        return False
    return max_abs_right < default_right - inch_to_emu(0.18)


def _title_align(shape, slide_width: int | None, explicit=None):
    align = _align_from_value(explicit)
    if align is not None:
        return align
    current = _current_align(shape)
    if slide_width:
        cx_ratio = (shape.left + shape.width / 2) / slide_width
        width_ratio = shape.width / slide_width
        if 0.42 <= cx_ratio <= 0.68 and width_ratio >= 0.32:
            return PP_ALIGN.CENTER
    return current or PP_ALIGN.LEFT


def _fit_body(shape, paragraphs, *, font, color, slide_width=None,
              visual_box=None, max_abs_right=None, max_abs_bottom=None,
              strict_bottom=False):
    base = round(_current_pt(shape, 16))
    max_right = (slide_width - inch_to_emu(0.35)) if slide_width else None
    respect_bounds = False
    measure_width_emu = None
    measure_height_emu = None
    max_measure_width_emu = None
    max_height_emu = None
    resize_width = True
    strict_height = bool(strict_bottom and visual_box and max_abs_bottom is not None)
    if visual_box and (max_abs_right is not None or strict_height):
        if max_abs_right is not None:
            max_right = min(max_right or max_abs_right, _local_limit_right(shape, visual_box, max_abs_right))
        respect_bounds = True
        measure_width_emu = visual_box[2]
        measure_height_emu = visual_box[3]
        if max_abs_right is not None:
            max_measure_width_emu = max(1, max_abs_right - visual_box[0])
        if strict_height:
            available_height = max_abs_bottom - visual_box[1]
            if available_height > 0:
                measure_height_emu = max(1, available_height)
                local_limit = _local_limit_bottom(shape, visual_box, max_abs_bottom)
                if local_limit and local_limit > shape.top:
                    max_height_emu = max(1, local_limit - shape.top)
        # Group children can use a local coordinate system whose dimensions do
        # not match the rendered slide-level box. Measure against the visual
        # box, but do not write that width back into the group-local geometry.
        if abs(shape.width - visual_box[2]) > max(1, visual_box[2] * 0.05):
            resize_width = False
    fit_body(
        shape, paragraphs,
        base_size_pt=base, min_size_pt=max(9, int(base * 0.6)),
        font_name=font, color=color, align=_current_align(shape) or PP_ALIGN.LEFT,
        max_right=max_right,
        respect_bounds=respect_bounds,
        measure_width_emu=measure_width_emu,
        measure_height_emu=measure_height_emu,
        max_measure_width_emu=max_measure_width_emu,
        resize_width=resize_width,
        strict_height=strict_height,
        max_height_emu=max_height_emu,
    )


def _slot_color(style: dict | None, fallback=None):
    if style and style.get("color"):
        return style["color"]
    return fallback


def _style_font(style: dict | None, fallback=None):
    if style and style.get("font_name"):
        return style["font_name"]
    return fallback


def _style_bold(style: dict | None, fallback=None):
    if style and style.get("bold") is not None:
        return bool(style["bold"])
    return fallback


def _split_lines(value) -> list[str]:
    return [line.strip() for line in str(value or "").splitlines() if line.strip()]


def _parallel_body_lines(value) -> list[str]:
    lines = _split_lines(value)
    if len(lines) != 1:
        return lines
    text = lines[0]
    compact = _norm_text(text)
    if len(compact) < 24:
        return lines
    midpoint = len(text) // 2
    candidates = []
    for token in ("；", "。", "，", "、", ";", ",", " "):
        start = 0
        while True:
            idx = text.find(token, start)
            if idx < 0:
                break
            split = idx + len(token)
            left = text[:split].strip()
            right = text[split:].strip()
            if len(_norm_text(left)) >= 8 and len(_norm_text(right)) >= 8:
                candidates.append((abs(split - midpoint), left, right))
            start = split
    if candidates:
        _distance, left, right = min(candidates, key=lambda item: item[0])
        return [left, right]
    return [text[:midpoint].strip(), text[midpoint:].strip()]


def _first_nonempty_paragraph_text(shape) -> str:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return ""
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            return text
    return ""


def _matches_first_template_paragraph(shape, value) -> bool:
    first = _first_nonempty_paragraph_text(shape)
    return bool(first) and _norm_text(first) == _norm_text(str(value or ""))


def _paragraph_has_style_profile(shape) -> bool:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return False
    paragraphs = list(shape.text_frame.paragraphs)
    nonempty = [p for p in paragraphs if p.text.strip()]
    if len(paragraphs) <= 1 or len(nonempty) <= 1:
        return False
    signatures = []
    for para in nonempty:
        run = next((r for r in para.runs if r.text.strip()), para.runs[0] if para.runs else None)
        size = run.font.size.pt if run is not None and run.font.size is not None else None
        font = run.font.name if run is not None else None
        bold = run.font.bold if run is not None else None
        color = _readable_run_color(run) if run is not None else None
        signatures.append((size, font, bold, color, para.alignment))
    return len(set(signatures)) > 1


def _readable_run_color(run):
    if run is None:
        return None
    try:
        color = run.font.color
        if color is None or color.type is None:
            return None
        try:
            if color.rgb is not None:
                return ("rgb", str(color.rgb))
        except Exception:
            pass
        try:
            if color.theme_color is not None:
                return ("theme", color.theme_color.name)
        except Exception:
            pass
    except Exception:
        pass
    return None


def _replace_paragraph_text(shape, values) -> None:
    """Replace existing paragraphs in-place, preserving per-paragraph styling."""
    text_frame = shape.text_frame
    values = [str(v) for v in _as_list(values)]
    while len(text_frame.paragraphs) < len(values):
        text_frame.add_paragraph()
    for i, value in enumerate(values):
        para = text_frame.paragraphs[i]
        if para.runs:
            para.runs[0].text = value
            for run in para.runs[1:]:
                run.text = ""
        else:
            para.text = value
    for para in text_frame.paragraphs[len(values):]:
        for run in para.runs:
            run.text = ""


def _set_paragraph_value(para, value: str) -> None:
    if para.runs:
        para.runs[0].text = value
        for run in para.runs[1:]:
            run.text = ""
    else:
        para.add_run().text = value


def _replace_indexed_paragraphs(shape, values_by_index: dict[int, str],
                                managed_indices: set[int]) -> None:
    text_frame = shape.text_frame
    for index, para in enumerate(text_frame.paragraphs):
        if index in values_by_index:
            _set_paragraph_value(para, values_by_index[index])
        elif index in managed_indices:
            for run in para.runs:
                run.text = ""


def _looks_numberish(text: str) -> bool:
    return bool(NUMBER_ONLY_RE.match(text or ""))


def _split_leading_number(text: str) -> tuple[str | None, str]:
    match = LEADING_NUMBER_RE.match(text or "")
    if not match:
        return None, text.strip()
    return match.group(1), match.group(2).strip()


def _bodyish_paragraph(text: str) -> bool:
    compact = _norm_text(text)
    return bool(BODYISH_TEXT_RE.search(text or "")) or len(compact) >= 22


def _replace_structured_card_text(shape, content: dict, *, preserve_number: bool = True) -> bool:
    """Fill a multi-paragraph card text box without flattening its styles.

    This handles common template boxes shaped as number/title/subtitle or
    number/title/body. Paragraphs keep their original run font, size, color,
    and alignment; only text changes.
    """
    if shape is None or not getattr(shape, "has_text_frame", False):
        return False
    paragraphs = list(shape.text_frame.paragraphs)
    nonempty = [(i, p.text.strip()) for i, p in enumerate(paragraphs) if p.text.strip()]
    if len(nonempty) < 2:
        return False

    first_idx, first_text = nonempty[0]
    has_number = _looks_numberish(first_text)
    title_lines = _split_lines(content.get("title"))
    body_lines = [str(v).strip() for v in _as_list(content.get("body")) if str(v).strip()]
    explicit_number = content.get("number")

    values: dict[int, str] = {}
    managed = {i for i, _text in nonempty}

    if has_number:
        number_value = str(explicit_number).strip() if explicit_number is not None else None
        if number_value is None and title_lines:
            parsed_number, remainder = _split_leading_number(title_lines[0])
            if parsed_number:
                number_value = parsed_number
                title_lines[0] = remainder
        if number_value is not None:
            values[first_idx] = number_value
        elif preserve_number:
            values[first_idx] = first_text
        title_candidates = [line for line in title_lines if line]
        title_idx = nonempty[1][0] if len(nonempty) > 1 else first_idx
        trailing = nonempty[2:]
    else:
        title_candidates = [line for line in title_lines if line]
        title_idx = first_idx
        trailing = nonempty[1:]

    if title_candidates:
        values[title_idx] = title_candidates[0]
    else:
        managed.discard(title_idx)

    trailing_indices = [i for i, _text in trailing]
    trailing_texts = [text for _i, text in trailing]
    has_body_region = any(_bodyish_paragraph(text) for text in trailing_texts)
    remainder_title = title_candidates[1:]
    if body_lines:
        detail_lines = body_lines
        detail_indices = trailing_indices
    elif has_body_region:
        detail_lines = remainder_title
        detail_indices = trailing_indices
    else:
        detail_lines = remainder_title
        detail_indices = trailing_indices

    for index, value in zip(detail_indices, detail_lines):
        values[index] = value
    if len(detail_lines) > len(detail_indices) and detail_indices:
        last = detail_indices[-1]
        extra = " ".join(detail_lines[len(detail_indices):])
        values[last] = (values.get(last, "") + " " + extra).strip()

    _replace_indexed_paragraphs(shape, values, managed)
    return bool(values)


def _replace_body_preserving_template(shape, value) -> bool:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return False
    lines = _split_lines(value)
    if not lines:
        return False
    paragraphs = list(shape.text_frame.paragraphs)
    nonempty = [i for i, p in enumerate(paragraphs) if p.text.strip()]
    if len(nonempty) <= 1:
        return False
    values = {}
    for index, line in zip(nonempty, lines):
        values[index] = line
    if len(lines) > len(nonempty):
        last = nonempty[-1]
        values[last] = (values[last] + " " + " ".join(lines[len(nonempty):])).strip()
    _replace_indexed_paragraphs(shape, values, set(nonempty))
    return True


def _normalize_text_font_size(shapes) -> None:
    """Make parallel text slots use the same smallest fitted font size."""
    text_shapes = [s for s in shapes if s is not None and getattr(s, "has_text_frame", False)]
    sizes = []
    for shape in text_shapes:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None and run.text.strip():
                    sizes.append(run.font.size.pt)
    if len(sizes) < 2:
        return
    target = min(sizes)
    for shape in text_shapes:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    run.font.size = Pt(target)


def _card_body_text(value) -> str:
    parts = [str(v).strip() for v in _as_list(value) if str(v).strip()]
    return "；".join(parts)


def _same_shape(a, b) -> bool:
    if a is None or b is None:
        return False
    if a is b:
        return True
    return getattr(a, "_element", None) is getattr(b, "_element", None)


def _has_icon_frame(slot: dict) -> bool:
    return bool(slot.get("icon_overlays"))


def _is_template_icon_frame(shape) -> bool:
    if shape is None:
        return False
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            return False
        if shape.fill.type == MSO_FILL.PICTURE:
            return False
        kind = getattr(shape.auto_shape_type, "name", str(shape.auto_shape_type))
        if kind not in {"OVAL", "ROUNDED_RECTANGLE", "RECTANGLE"}:
            return False
        aspect = shape.width / shape.height if shape.height else 0
        return 0.75 <= aspect <= 1.35
    except Exception:
        return False


def _shape_box(shape) -> tuple[int, int, int, int] | None:
    if shape is None:
        return None
    return (shape.left, shape.top, shape.width, shape.height)


def _geometry_box(geometry: dict | None) -> tuple[int, int, int, int] | None:
    if not geometry:
        return None
    try:
        return (
            int(geometry["left"]),
            int(geometry["top"]),
            int(geometry["width"]),
            int(geometry["height"]),
        )
    except Exception:
        return None


def _path_tuple(path) -> tuple:
    if path is None:
        return ()
    return tuple(path)


def _overlap(a_start: int, a_size: int, b_start: int, b_size: int) -> int:
    return max(0, min(a_start + a_size, b_start + b_size) - max(a_start, b_start))


def _active_text_layout_boxes(slide, fill_plan: dict, page: dict) -> list[dict]:
    """Slide-level text boxes and visual obstacles that constrain fitting."""
    boxes: list[dict] = []

    def add(slot, kind: str):
        if not slot or not slot.get("path"):
            return
        try:
            boxes.append({
                "path": _path_tuple(slot["path"]),
                "kind": kind,
                "box": absolute_box_for_path(slide, slot["path"]),
            })
        except Exception:
            return

    def add_box(path, kind: str, box):
        if not box:
            return
        boxes.append({
            "path": _path_tuple(path),
            "kind": kind,
            "box": box,
        })

    if page.get("title"):
        add(fill_plan.get("title"), "title")
    if page.get("subtitle"):
        add(fill_plan.get("subtitle"), "subtitle")
    if page.get("section_number"):
        add(fill_plan.get("section_number"), "section_number")
    if page.get("footer"):
        add(fill_plan.get("footer"), "footer")

    labels = _as_list(page.get("labels")) if "labels" in page else []
    for slot, value in zip(fill_plan.get("labels", []), labels):
        if value:
            add(slot, "label")

    body_values = [p for p in _as_list(page.get("body")) if str(p).strip()] if "body" in page else []
    for slot, value in zip(fill_plan.get("body", []), body_values):
        if value:
            add(slot, "body")

    for card_slot, card_content in zip(fill_plan.get("cards", []), page.get("cards") or []):
        if card_content.get("title") and card_slot.get("title"):
            add({"path": card_slot["title"]}, "card_title")
        if card_content.get("body") and card_slot.get("body"):
            add({"path": card_slot["body"]}, "card_body")
        if card_content.get("number") and card_slot.get("number"):
            add({"path": card_slot["number"]}, "card_number")

    for image_slot in fill_plan.get("images", []):
        try:
            add_box(image_slot.get("path"), "image", absolute_box_for_path(slide, image_slot["path"]))
        except Exception:
            continue

    for obstacle in fill_plan.get("visual_obstacles", []):
        box = None
        try:
            if obstacle.get("path"):
                box = absolute_box_for_path(slide, obstacle["path"])
        except Exception:
            box = None
        box = box or _geometry_box(obstacle.get("geometry"))
        add_box(obstacle.get("path"), f"visual:{obstacle.get('kind') or 'shape'}", box)
    return boxes


def _layout_limits(own_box, layout_boxes: list[dict], *, own_path=None,
                   slide_width=None, slide_height=None,
                   container_box=None) -> dict:
    """Find nearest right/below text neighbors and return slide-level limits."""
    if not own_box:
        return {}
    left, top, width, height = own_box
    gutter = inch_to_emu(0.08)
    limits: dict = {}
    if slide_width:
        limits["max_abs_right"] = slide_width - inch_to_emu(0.35)
    if slide_height:
        limits["max_abs_bottom"] = slide_height - inch_to_emu(0.25)
    if container_box:
        c_left, c_top, c_width, c_height = container_box
        if left >= c_left - gutter and top < c_top + c_height and top + height > c_top:
            candidate = c_left + c_width - gutter
            if candidate > left + inch_to_emu(0.25):
                limits["max_abs_right"] = min(limits.get("max_abs_right", candidate), candidate)
        if top >= c_top - gutter and left < c_left + c_width and left + width > c_left:
            candidate = c_top + c_height - gutter
            if candidate > top + inch_to_emu(0.18):
                if candidate < limits.get("max_abs_bottom", candidate + 1):
                    limits["max_abs_bottom"] = candidate
                    limits["bottom_constraint_kind"] = "container"

    own_path = _path_tuple(own_path)
    for item in layout_boxes:
        if own_path and item.get("path") == own_path:
            continue
        other = item.get("box")
        if not other:
            continue
        o_left, o_top, o_width, o_height = other
        v_ov = _overlap(top, height, o_top, o_height)
        h_ov = _overlap(left, width, o_left, o_width)
        # Right-side neighbor: same row/vertical band and visually to the right.
        #
        # Wide title placeholders often already extend into a right-side image
        # or decorative panel. In that case the neighbor's left edge can sit
        # inside the original title box, so checking only "to the right of the
        # current right edge" misses the real layout boundary.
        same_vertical_band = v_ov >= min(height, o_height) * 0.22
        right_edge_after_own_left = o_left > left + inch_to_emu(0.12)
        other_center_right_of_own = (o_left + o_width / 2) > (left + width / 2) + inch_to_emu(0.05)
        if (
            same_vertical_band
            and right_edge_after_own_left
            and other_center_right_of_own
        ):
            candidate = o_left - gutter
            if candidate > left + inch_to_emu(0.25):
                limits["max_abs_right"] = min(
                    limits.get("max_abs_right", candidate),
                    candidate,
                )
                limits["right_constraint_kind"] = item.get("kind")
        # Below neighbor: overlapping column and lower on the page.
        if (
            h_ov >= min(width, o_width) * 0.22
            and o_top > top + height * 0.12
        ):
            candidate = o_top - gutter
            if candidate > top + inch_to_emu(0.18):
                if candidate < limits.get("max_abs_bottom", candidate + 1):
                    limits["max_abs_bottom"] = candidate
                    limits["bottom_constraint_kind"] = item.get("kind")
    return limits


def _label_row_limits(slide, label_slots: list[dict],
                      labels_content: list, *, page_role=None) -> dict[tuple, int]:
    """Constrain parallel label rows to their template column width."""
    if len(labels_content) < 3:
        return {}
    rows: dict[int, list[tuple[tuple, tuple[int, int, int, int]]]] = {}
    for slot, value in zip(label_slots, labels_content):
        if not value or not slot or not slot.get("path"):
            continue
        try:
            box = absolute_box_for_path(slide, slot["path"])
        except Exception:
            box = None
        if not box:
            continue
        row = round(box[1] / inch_to_emu(0.25))
        rows.setdefault(row, []).append((_path_tuple(slot["path"]), box))

    limits: dict[tuple, int] = {}
    for items in rows.values():
        if len(items) < 3:
            continue
        widths = sorted(box[2] for _path, box in items)
        median_width = widths[len(widths) // 2]
        for path, box in items:
            limits[path] = box[0] + median_width
    return limits


def _normalize_parallel_label_rows(slide, label_shapes, label_slots: list[dict],
                                   filled_count: int) -> None:
    groups: dict[tuple, list] = {}
    for shape, slot in zip(label_shapes[:filled_count], label_slots[:filled_count]):
        if shape is None or not slot or not slot.get("path"):
            continue
        sample = str(slot.get("sample") or "").strip()
        current = shape.text_frame.text.strip() if getattr(shape, "has_text_frame", False) else ""
        if _looks_numberish(sample) or _looks_numberish(current):
            continue
        style = slot.get("style") or {}
        try:
            box = absolute_box_for_path(slide, slot["path"])
        except Exception:
            box = None
        if not box:
            continue
        key = (
            round(float(style.get("font_pt") or 0), 1),
            bool(style.get("bold")),
            str(style.get("align") or ""),
            round(box[3] / inch_to_emu(0.12)),
        )
        groups.setdefault(key, []).append(shape)
    for shapes in groups.values():
        if len(shapes) >= 2:
            _normalize_text_font_size(shapes)


def _local_limit_right(shape, visual_box, max_abs_right):
    if shape is None or not visual_box or max_abs_right is None:
        return max_abs_right
    abs_left, _abs_top, abs_width, _abs_height = visual_box
    scale_x = abs_width / shape.width if shape.width else 1.0
    if not scale_x:
        return max_abs_right
    return int(shape.left + (max_abs_right - abs_left) / scale_x)


def _local_limit_bottom(shape, visual_box, max_abs_bottom):
    if shape is None or not visual_box or max_abs_bottom is None:
        return max_abs_bottom
    _abs_left, abs_top, _abs_width, abs_height = visual_box
    scale_y = abs_height / shape.height if shape.height else 1.0
    if not scale_y:
        return max_abs_bottom
    return int(shape.top + (max_abs_bottom - abs_top) / scale_y)


def _similar_dimension(values: list[int], tolerance: float = 0.22) -> bool:
    if len(values) < 2:
        return False
    lo, hi = min(values), max(values)
    return hi > 0 and (hi - lo) <= hi * tolerance


def _same_row(boxes: list[tuple[int, int, int, int]]) -> bool:
    if len(boxes) < 2:
        return False
    centers_y = [top + height // 2 for _left, top, _width, height in boxes]
    heights = [height for _left, _top, _width, height in boxes]
    return max(centers_y) - min(centers_y) <= max(heights) * 0.35


def _same_row_or_column(boxes: list[tuple[int, int, int, int]]) -> bool:
    if len(boxes) < 2:
        return False
    centers_x = [left + width // 2 for left, _top, width, _height in boxes]
    widths = [width for _left, _top, width, _height in boxes]
    same_column = max(centers_x) - min(centers_x) <= max(widths) * 0.35
    return _same_row(boxes) or same_column


def _normalize_parallel_filled_text_boxes(
    items: list[tuple[object, tuple[int, int, int, int] | None]],
    *,
    rows_only: bool = False,
) -> None:
    """Normalize fitted font size for same-row/same-geometry text peers."""
    clean_items: list[tuple[object, tuple[int, int, int, int]]] = []
    for shape, box in items:
        if shape is None or not getattr(shape, "has_text_frame", False) or not box:
            continue
        current = shape.text_frame.text.strip()
        if not current or _looks_numberish(current):
            continue
        clean_items.append((shape, box))

    groups: list[list[tuple[object, tuple[int, int, int, int]]]] = []
    for item in clean_items:
        _shape, box = item
        placed = False
        for group in groups:
            boxes = [b for _s, b in group] + [box]
            same_peer_band = _same_row(boxes) if rows_only else _same_row_or_column(boxes)
            if (
                same_peer_band
                and _similar_dimension([b[2] for b in boxes])
                and _similar_dimension([b[3] for b in boxes])
            ):
                group.append(item)
                placed = True
                break
        if not placed:
            groups.append([item])

    for group in groups:
        if len(group) >= 2:
            _normalize_text_font_size([shape for shape, _box in group])


def _normalize_parallel_filled_text_peers(slide, entries: list[tuple[object, dict | None]],
                                          initial_boxes: dict[tuple, tuple[int, int, int, int]] | None = None,
                                          *, rows_only: bool = False) -> None:
    """Normalize same-size peer text boxes even when parser split their roles."""
    items: list[tuple[object, tuple[int, int, int, int]]] = []
    for shape, slot in entries:
        if shape is None or not getattr(shape, "has_text_frame", False):
            continue
        current = shape.text_frame.text.strip()
        if not current or _looks_numberish(current):
            continue
        sample = str((slot or {}).get("sample") or "").strip()
        if _looks_numberish(sample):
            continue
        if not slot or not slot.get("path"):
            continue
        path = _path_tuple(slot["path"])
        box = (initial_boxes or {}).get(path)
        if box is None:
            try:
                box = absolute_box_for_path(slide, slot["path"])
            except Exception:
                box = None
        if not box:
            continue
        items.append((shape, box))

    _normalize_parallel_filled_text_boxes(items, rows_only=rows_only)


def _use_bounded_card_titles(resolved_cards, cards_content, *, page_role=None) -> bool:
    """Repeated title-only card lists are layout-constrained, not free labels."""
    pairs = [
        (slot, content)
        for slot, content in zip(resolved_cards, cards_content)
        if slot.get("title") is not None and content.get("title")
    ]
    if len(pairs) < 2:
        return False
    if any(slot.get("body") is not None or content.get("body") for slot, content in pairs):
        return False
    boxes = [slot.get("title_box") or _shape_box(slot["title"]) for slot, _content in pairs]
    boxes = [box for box in boxes if box is not None]
    if len(boxes) < 2:
        return False
    widths = [box[2] for box in boxes]
    heights = [box[3] for box in boxes]
    repeated_layout = (
        _similar_dimension(widths)
        and _similar_dimension(heights)
        and _same_row_or_column(boxes)
    )
    return page_role == "contents" or repeated_layout


def _replace_card_icon(slide, slot: dict, png) -> None:
    for overlay in slot.get("icon_overlays", []):
        if overlay is not None and not _same_shape(overlay, slot.get("icon")):
            try:
                delete_shape(overlay)
            except Exception:
                hide_shape_visual(overlay)
    box = slot.get("icon_box") or _shape_box(slot.get("icon"))
    if box is None:
        return
    if _has_icon_frame(slot) or _is_template_icon_frame(slot.get("icon")):
        _clear_text(slot.get("icon"))
        add_centered_picture(
            slide, png,
            box[0], box[1], box[2], box[3],
            padding_ratio=0.28,
        )
    else:
        try:
            delete_shape(slot["icon"])
        except Exception:
            hide_shape_visual(slot["icon"])
        add_centered_picture(
            slide, png,
            box[0], box[1], box[2], box[3],
            padding_ratio=0.2,
        )


def _clear_card_residual_text(slot: dict) -> int:
    cleared = 0
    protected = [
        shape for shape in (slot.get("title"), slot.get("body"), slot.get("number"))
        if shape is not None
    ]
    for shape in slot.get("residual_text") or []:
        if shape is None or any(_same_shape(shape, p) for p in protected):
            continue
        if _clear_text(shape):
            cleared += 1
    return cleared


def _icon_color_for_slot(slot: dict, brand: dict) -> tuple[int, int, int] | None:
    explicit = brand.get("icon_color")
    if explicit:
        return _rgb_tuple(explicit)
    return _rgb_tuple((slot.get("icon_style") or {}).get("color"))


def _fill_cards(slide, resolved_cards, cards_content, *, asset_dir, brand, warnings,
                slide_width=None, slide_height=None, page_role=None, layout_boxes=None):
    accent = _to_rgb(brand.get("accent"))
    title_font = brand.get("title_font")
    body_font = brand.get("body_font")
    filled = min(len(resolved_cards), len(cards_content))
    residual_cleared = 0
    bounded_card_titles = _use_bounded_card_titles(
        resolved_cards[:filled],
        cards_content[:filled],
        page_role=page_role,
    )
    filled_card_title_boxes: list[tuple[object, tuple[int, int, int, int] | None]] = []

    for i in range(filled):
        slot, content = resolved_cards[i], cards_content[i]
        combined_body = _card_body_text(content.get("body"))
        shared_textbox = (
            slot.get("title_body_mode") == "shared_textbox"
            or (not slot.get("title_body_mode") and (not slot["body"] or _same_shape(slot["body"], slot["title"])))
        )
        if slot["title"] and content.get("title"):
            title_style = slot.get("title_style") or {}
            title_font = _style_font(title_style, brand.get("title_font"))
            title_bold = _style_bold(title_style, None)
            title_color = _slot_color(title_style, accent)
            if shared_textbox and _replace_structured_card_text(slot["title"], content):
                pass
            elif _paragraph_has_style_profile(slot["title"]) and _replace_structured_card_text(slot["title"], content):
                pass
            elif combined_body and shared_textbox:
                _replace_paragraph_text(slot["title"], [content["title"], combined_body])
            else:
                title_limits = _layout_limits(
                    slot.get("title_box"),
                    layout_boxes or [],
                    own_path=None,
                    slide_width=slide_width,
                    slide_height=slide_height,
                    container_box=slot.get("card_box"),
                )
                if bounded_card_titles:
                    _fit_bounded_title(
                        slot["title"],
                        content["title"],
                        font=title_font,
                        color=title_color,
                        visual_box=slot.get("title_box"),
                        bold=title_bold,
                    )
                else:
                    max_right = _local_limit_right(
                        slot["title"],
                        slot.get("title_box"),
                        title_limits.get("max_abs_right"),
                    )
                    max_bottom = _local_limit_bottom(
                        slot["title"],
                        slot.get("title_box"),
                        title_limits.get("max_abs_bottom"),
                    )
                    title_default_box = slot.get("card_box") or slot.get("title_box")
                    title_default_right = (
                        title_default_box[0] + title_default_box[2]
                        if title_default_box else None
                    )
                    title_right_constrained = _right_limit_is_constrained(
                        title_limits,
                        title_default_right,
                    )
                    _fit_title(
                        slot["title"],
                        content["title"],
                        font=title_font,
                        color=title_color,
                        max_right=max_right,
                        max_bottom=max_bottom,
                        allow_break=True,
                        prefer_break=title_right_constrained,
                        respect_bounds=bool(title_limits),
                        bold=title_bold,
                    )
            filled_card_title_boxes.append((slot["title"], slot.get("title_box")))
        elif slot["title"]:
            _clear_if_stock(slot["title"])
        if slot["number"] and content.get("number") is not None:
            replace_first_run_text(slot["number"], str(content["number"]))
        body_content = _as_list(content.get("body")) if content.get("body") else []
        if slot["body"] and not slot["title"] and content.get("title"):
            body_content = [str(content["title"]), *body_content]
        if slot["body"] and not shared_textbox and body_content:
            body_limits = _layout_limits(
                slot.get("body_box"),
                layout_boxes or [],
                own_path=None,
                slide_width=slide_width,
                slide_height=slide_height,
                container_box=slot.get("card_box"),
            )
            _fit_body(
                slot["body"],
                body_content,
                font=_style_font(slot.get("body_style"), body_font),
                color=_slot_color(slot.get("body_style")),
                slide_width=slide_width,
                visual_box=slot.get("body_box"),
                max_abs_right=body_limits.get("max_abs_right"),
                max_abs_bottom=body_limits.get("max_abs_bottom"),
                strict_bottom=bool(body_limits.get("bottom_constraint_kind")),
            )
        elif slot["body"]:
            _clear_if_stock(slot["body"])
        if content.get("icon"):
            if slot["icon"] is not None:
                color = _icon_color_for_slot(slot, brand)
                png = get_icon(content["icon"], asset_dir, color=color)
                _replace_card_icon(slide, slot, png)
            else:
                warnings.append(f"card {i}: icon '{content['icon']}' requested but slot has no icon")
        residual_cleared += _clear_card_residual_text(slot)

    if bounded_card_titles and filled_card_title_boxes:
        _normalize_parallel_filled_text_boxes(
            filled_card_title_boxes,
            rows_only=(page_role == "contents"),
        )

    # clear leftover placeholders in card slots that received no content
    for slot in resolved_cards[filled:]:
        for key in ("title", "body", "number"):
            _clear_if_stock(slot[key])

    if len(cards_content) > len(resolved_cards):
        warnings.append(f"{len(cards_content) - len(resolved_cards)} extra cards dropped "
                        f"(template has {len(resolved_cards)} card slots)")
    return filled, residual_cleared


def _fill_body_slots(resolved_body, paragraphs, *, font, warnings, slots=None,
                     slide_width=None, slide_height=None, slide=None, layout_boxes=None,
                     page_role=None):
    shapes = [s for s in resolved_body if s is not None]
    slots = slots or []
    if not shapes:
        if paragraphs:
            warnings.append("body content given but slide has no body slot")
        return 0
    # filter out empty paragraphs
    paragraphs = [p for p in paragraphs if p.strip()]
    if not paragraphs:
        return 0
    if len(shapes) == 1:
        slot = slots[0] if slots else None
        if _replace_body_preserving_template(shapes[0], "\n".join(paragraphs)):
            return 1
        own_box = absolute_box_for_path(slide, slot["path"]) if slide is not None and slot else None
        limits = _layout_limits(
            own_box,
            layout_boxes or [],
            own_path=slot.get("path") if slot else None,
            slide_width=slide_width,
            slide_height=slide_height,
        )
        _fit_body(
            shapes[0],
            paragraphs,
            font=_style_font((slot or {}).get("style"), font),
            color=_slot_color((slot or {}).get("style")),
            slide_width=slide_width,
            visual_box=own_box,
            max_abs_right=limits.get("max_abs_right"),
            max_abs_bottom=limits.get("max_abs_bottom"),
            strict_bottom=bool(limits.get("bottom_constraint_kind")),
        )
        return 1
    count = min(len(shapes), len(paragraphs))
    for i in range(count):
        slot = slots[i] if i < len(slots) else None
        if _replace_body_preserving_template(shapes[i], paragraphs[i]):
            continue
        own_box = absolute_box_for_path(slide, slot["path"]) if slide is not None and slot else None
        limits = _layout_limits(
            own_box,
            layout_boxes or [],
            own_path=slot.get("path") if slot else None,
            slide_width=slide_width,
            slide_height=slide_height,
        )
        _fit_body(
            shapes[i],
            _parallel_body_lines(paragraphs[i]),
            font=_style_font((slot or {}).get("style"), font),
            color=_slot_color((slot or {}).get("style")),
            slide_width=slide_width,
            visual_box=own_box,
            max_abs_right=(own_box[0] + own_box[2]) if own_box else limits.get("max_abs_right"),
            max_abs_bottom=limits.get("max_abs_bottom"),
            strict_bottom=bool(limits.get("bottom_constraint_kind")),
        )
    if len(paragraphs) > len(shapes):
        warnings.append(f"{len(paragraphs) - len(shapes)} extra body paragraphs dropped")
    if page_role == "contents" and count >= 2:
        _normalize_parallel_filled_text_peers(
            slide,
            [
                (shapes[i], slots[i] if i < len(slots) else None)
                for i in range(count)
            ],
            rows_only=True,
        )
    return count


def _image_item_path(item) -> str | None:
    if item is None:
        return None
    if isinstance(item, str):
        return item.strip() or None
    if isinstance(item, dict):
        path = str(item.get("path") or "").strip()
        if path:
            return path
        return None
    return str(item).strip() or None


def _image_slot_shape(slot):
    if isinstance(slot, dict):
        return slot.get("shape")
    return slot


def _image_slot_box(slot):
    if isinstance(slot, dict):
        return slot.get("box")
    shape = _image_slot_shape(slot)
    if shape is None:
        return None
    return (shape.left, shape.top, shape.width, shape.height)


def _fill_images(slide, resolved_images, image_items, *, warnings):
    slots = [s for s in resolved_images if _image_slot_shape(s) is not None]
    count = 0
    slot_count = min(len(slots), len(image_items))
    for i in range(slot_count):
        image_path = _image_item_path(image_items[i])
        if not image_path:
            continue
        if not Path(image_path).exists():
            warnings.append(f"image {i}: file not found, kept original ({image_path})")
            continue
        slot = slots[i]
        shape = _image_slot_shape(slot)
        box = _image_slot_box(slot)
        if not box:
            warnings.append(f"image {i}: slot geometry missing, kept original")
            continue
        left, top, width, height = box
        target_ratio = (width / height) if height else 1.0
        cropped = crop_to_fill(image_path, f"{image_path}.crop_{i}.png",
                               target_ratio=target_ratio)
        slide.shapes.add_picture(str(cropped), left, top, width=width, height=height)
        delete_shape(shape)
        count += 1
    extra = [
        item for item in image_items[len(slots):]
        if _image_item_path(item)
    ]
    if extra:
        warnings.append(f"{len(extra)} extra images dropped")
    return count


# ---------------------------------------------------------------------------
# top-level: apply one page
# ---------------------------------------------------------------------------
def apply_page(slide, fill_plan: dict, page: dict, *, asset_dir,
               brand: dict | None = None, slide_width: int | None = None,
               slide_height: int | None = None,
               clear_unfilled: bool = True) -> dict:
    """Fill `slide` from `page` content against its `fill_plan`. Returns a report."""
    brand = brand or {}
    warnings: list[str] = []
    report = {"slide": page.get("slide"), "role": page.get("role"), "filled": {}, "warnings": warnings}

    res = resolve_fill_plan(slide, fill_plan)
    layout_boxes = _active_text_layout_boxes(slide, fill_plan, page)
    initial_boxes = {
        item["path"]: item["box"]
        for item in layout_boxes
        if item.get("path") and item.get("box")
    }
    accent = _to_rgb(brand.get("accent"))
    title_font = brand.get("title_font")
    body_font = brand.get("body_font")
    title_min_left = inch_to_emu(0.4)
    title_max_right = (slide_width - inch_to_emu(0.4)) if slide_width else None

    # title
    if page.get("title") and res["title"] is not None:
        title_text = str(page["title"])
        title_lines = [line for line in title_text.splitlines() if line.strip()]
        title_style = (fill_plan.get("title") or {}).get("style") or {}
        title_color = _slot_color(title_style, accent)
        title_font_for_slot = _style_font(title_style, title_font)
        title_bold = _style_bold(title_style, None)
        title_align = _title_align(res["title"], slide_width, page.get("title_align"))
        title_slot = fill_plan.get("title") or {}
        title_box = absolute_box_for_path(slide, title_slot["path"]) if title_slot.get("path") else None
        title_limits = _layout_limits(
            title_box,
            layout_boxes,
            own_path=title_slot.get("path"),
            slide_width=slide_width,
            slide_height=slide_height,
        )
        title_right_constrained = _right_limit_is_constrained(title_limits, title_max_right)
        local_max_right = _local_limit_right(
            res["title"],
            title_box,
            title_limits.get("max_abs_right"),
        )
        local_max_bottom = _local_limit_bottom(
            res["title"],
            title_box,
            title_limits.get("max_abs_bottom"),
        )
        if _matches_first_template_paragraph(res["title"], title_text):
            pass
        elif len(title_lines) >= 2 and len(res["title"].text_frame.paragraphs) >= len(title_lines):
            _replace_paragraph_text(res["title"], title_lines)
        else:
            _fit_title(res["title"], page["title"], font=title_font_for_slot, color=title_color,
                       align=title_align,
                       min_left=title_min_left,
                       max_right=min(title_max_right, local_max_right) if title_max_right and local_max_right else (local_max_right or title_max_right),
                       max_bottom=local_max_bottom,
                       allow_break=_title_allows_break(
                           res["title"],
                           slide_width,
                           right_constrained=title_right_constrained,
                       ),
                       prefer_break=title_right_constrained,
                       respect_bounds=bool(title_limits),
                       bold=title_bold,
                       min_size_ratio=0.45 if title_right_constrained else 0.6)
        report["filled"]["title"] = True
    elif clear_unfilled:
        _clear_if_unfilled(res["title"], fill_plan.get("title"))

    # subtitle (inherit color)
    if page.get("subtitle") and res["subtitle"] is not None:
        subtitle_style = (fill_plan.get("subtitle") or {}).get("style") or {}
        set_text(
            res["subtitle"],
            _as_list(page["subtitle"]),
            font_name=_style_font(subtitle_style, title_font),
            bold=_style_bold(subtitle_style),
            color=_slot_color(subtitle_style),
        )
        report["filled"]["subtitle"] = True
    elif clear_unfilled:
        _clear_if_unfilled(res["subtitle"], fill_plan.get("subtitle"))

    # section number — strong template element, keep run styling
    if page.get("section_number") and res["section_number"] is not None:
        replace_first_run_text(res["section_number"], str(page["section_number"]))
        report["filled"]["section_number"] = True

    # footer
    if res["footer"] is not None and "footer" in page and _has_content(page.get("footer")):
        footer_style = (fill_plan.get("footer") or {}).get("style") or {}
        set_text(
            res["footer"],
            _as_list(page["footer"]),
            font_name=_style_font(footer_style),
            bold=_style_bold(footer_style),
            color=_slot_color(footer_style),
        )
        report["filled"]["footer"] = True
    elif res["footer"] is not None:
        report["filled"]["footer"] = "preserved"

    # labels - treat like body slots but simpler (no autofit)
    label_slots = fill_plan.get("labels", [])
    if "labels" in page:
        labels_content = _as_list(page["labels"])
        label_row_limits = _label_row_limits(
            slide,
            label_slots,
            labels_content,
            page_role=page.get("role"),
        )
        n = 0
        for i, shape in enumerate(res["labels"]):
            slot = label_slots[i] if i < len(label_slots) else None
            if shape is not None and i < len(labels_content):
                content = labels_content[i]
                if content:
                    if (
                        slot
                        and slot.get("default_action") == "preserve"
                        and _norm_text(content) == _norm_text(slot.get("sample"))
                        and _norm_text(shape.text_frame.text) == _norm_text(slot.get("sample"))
                    ):
                        n += 1
                        continue
                    text_value = str(content)
                    if len(_norm_text(text_value)) >= 4:
                        label_box = absolute_box_for_path(slide, slot["path"]) if slot else None
                        label_limits = _layout_limits(
                            label_box,
                            layout_boxes,
                            own_path=slot.get("path") if slot else None,
                            slide_width=slide_width,
                            slide_height=slide_height,
                        )
                        row_limit = label_row_limits.get(_path_tuple(slot.get("path") if slot else None))
                        if row_limit is not None:
                            label_limits["max_abs_right"] = min(
                                label_limits.get("max_abs_right", row_limit),
                                row_limit,
                            )
                        _fit_label(
                            shape,
                            text_value,
                            font=_style_font((slot or {}).get("style"), title_font),
                            slide_width=slide_width,
                            color=_slot_color((slot or {}).get("style")),
                            visual_box=label_box,
                            max_abs_right=label_limits.get("max_abs_right"),
                            max_abs_bottom=label_limits.get("max_abs_bottom"),
                        )
                    else:
                        replace_first_run_text(shape, text_value)
                        _set_word_wrap(shape, False)
                    n += 1
                elif content == "":
                    _clear_text(shape)
                elif clear_unfilled:
                    _clear_if_unfilled(shape, slot)
            elif clear_unfilled:
                _clear_if_unfilled(shape, slot)
        if n > 0:
            _normalize_parallel_label_rows(slide, res["labels"], label_slots, n)
            peer_entries = []
            if page.get("title") and res["title"] is not None:
                peer_entries.append((res["title"], fill_plan.get("title")))
            for i, shape in enumerate(res["labels"][:n]):
                slot = label_slots[i] if i < len(label_slots) else None
                peer_entries.append((shape, slot))
            _normalize_parallel_filled_text_peers(
                slide,
                peer_entries,
                initial_boxes,
                rows_only=(page.get("role") == "contents"),
            )
            report["filled"]["labels"] = n
    elif clear_unfilled:
        for i, shape in enumerate(res["labels"]):
            slot = label_slots[i] if i < len(label_slots) else None
            _clear_if_unfilled(shape, slot)

    # body
    body_slots = fill_plan.get("body", [])
    if "body" in page:
        paragraphs = [p for p in _as_list(page.get("body")) if str(p).strip()]
        if paragraphs:
            n = _fill_body_slots(
                res["body"],
                paragraphs,
                font=body_font,
                warnings=warnings,
                slots=body_slots,
                slide_width=slide_width,
                slide_height=slide_height,
                slide=slide,
                layout_boxes=layout_boxes,
                page_role=page.get("role"),
            )
        else:
            n = 0
            if clear_unfilled:
                for shape in res["body"]:
                    _clear_text(shape)
        report["filled"]["body"] = n
    elif clear_unfilled:
        for i, shape in enumerate(res["body"]):
            slot = body_slots[i] if i < len(body_slots) else None
            _clear_if_unfilled(shape, slot)

    # cards
    if page.get("cards"):
        n, residual_cleared = _fill_cards(
            slide, res["cards"], page["cards"],
            asset_dir=asset_dir, brand=brand, warnings=warnings,
            slide_width=slide_width, slide_height=slide_height,
            page_role=page.get("role"), layout_boxes=layout_boxes,
        )
        report["filled"]["cards"] = n
        if residual_cleared:
            report["filled"]["card_residual_text_cleared"] = residual_cleared
    elif clear_unfilled:
        for slot in res["cards"]:
            for key in ("title", "body", "number"):
                _clear_if_stock(slot[key])

    # images
    if page.get("images"):
        n = _fill_images(slide, res["images"], page["images"], warnings=warnings)
        report["filled"]["images"] = n

    if fill_plan.get("charts") and not page.get("charts"):
        warnings.append("chart slot present but page has no chart data; choose another page or implement chart replacement")
    if fill_plan.get("tables") and not page.get("tables"):
        warnings.append("table slot present but page has no table data; choose another page or implement table replacement")

    # logo (centered, light padding)
    if page.get("logo") and res["logo"] is not None:
        replace_shape_with_picture(slide, res["logo"], page["logo"],
                                   centered=True, padding_ratio=0.08)
        report["filled"]["logo"] = True
    elif clear_unfilled:
        _clear_if_stock(res["logo"])

    if clear_unfilled:
        residual = _clear_stock_texts(slide.shapes)
        if residual:
            report["filled"]["cleared_placeholders"] = residual

    return report
