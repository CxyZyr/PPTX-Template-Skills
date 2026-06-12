"""Shape-tree walking and feature extraction for template parsing.

New module (the parsing side). Provides a recursive walk that yields stable
index paths, plus normalized geometry / text / type features that the role
classifier consumes. Everything is defensive: malformed shapes never raise.
"""
from __future__ import annotations

from collections import Counter

from .geometry import emu_to_inch

try:
    from pptx.oxml.ns import qn
except Exception:  # pragma: no cover
    qn = None

try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    _GROUP = MSO_SHAPE_TYPE.GROUP
except Exception:  # pragma: no cover
    _GROUP = None


def is_group(shape) -> bool:
    try:
        return shape.shape_type == _GROUP
    except Exception:
        return False


def iter_shapes(shapes, prefix: tuple = ()):
    """Yield (path, shape) for every shape, descending into groups.

    path is a tuple of indexes, e.g. (4,) for a top-level shape and (4, 1) for
    the 2nd child of group 4. z-order == iteration order within a container.
    """
    for index, shape in enumerate(shapes):
        path = prefix + (index,)
        yield path, shape
        if is_group(shape):
            yield from iter_shapes(shape.shapes, path)


def iter_top_level(shapes):
    """Yield (index, shape) for direct children only (no recursion)."""
    for index, shape in enumerate(shapes):
        yield index, shape


def get_shape_by_path(slide_or_group, path):
    """Resolve an index-path (e.g. [2, 1]) to the actual shape object.

    This is the inverse of the path emitted by `iter_shapes` and stored in the
    parsing spec: [2, 1] means slide.shapes[2].shapes[1]. Descends into groups
    by index. Returns None if any index is out of range or a non-group is asked
    to descend — generation code should treat None as "slot missing" and skip.
    """
    shapes = getattr(slide_or_group, "shapes", slide_or_group)
    shape = None
    for depth, index in enumerate(path):
        try:
            shape = shapes[index]
        except (IndexError, TypeError, KeyError):
            return None
        if depth < len(path) - 1:
            if not is_group(shape):
                return None
            shapes = shape.shapes
    return shape


def resolve_paths(slide, paths):
    """Resolve many index-paths at once, returning a dict {tuple(path): shape}.

    Resolve everything up front (one read pass) BEFORE mutating the slide: held
    shape references stay valid when siblings are deleted, but re-resolving by
    index after a delete would point at the wrong shape.
    """
    out = {}
    for path in paths:
        if path is None:
            continue
        out[tuple(path)] = get_shape_by_path(slide, path)
    return out


def _box_tuple(shape) -> tuple[int, int, int, int]:
    return (
        int(shape.left or 0),
        int(shape.top or 0),
        int(shape.width or 0),
        int(shape.height or 0),
    )


def _group_child_transform(shape) -> tuple[float, float, float, float] | None:
    """Return (scale_x, scale_y, child_off_x, child_off_y) for a group.

    Group children are stored in the group's child coordinate space. To place
    a replacement picture on the slide we need to map that child space back to
    slide coordinates instead of using python-pptx's local child left/top.
    """
    if qn is None or not is_group(shape):
        return None
    try:
        xfrm = shape._element.grpSpPr.xfrm  # noqa: SLF001
        ch_off = xfrm.find(qn("a:chOff"))
        ch_ext = xfrm.find(qn("a:chExt"))
        if ch_off is None or ch_ext is None:
            return None
        _, _, width, height = _box_tuple(shape)
        child_w = int(ch_ext.get("cx") or 0)
        child_h = int(ch_ext.get("cy") or 0)
        sx = width / child_w if child_w else 1.0
        sy = height / child_h if child_h else 1.0
        return sx, sy, int(ch_off.get("x") or 0), int(ch_off.get("y") or 0)
    except Exception:
        return None


def absolute_box_for_path(slide_or_group, path) -> tuple[int, int, int, int] | None:
    """Resolve `path` and return its slide-level (left, top, width, height).

    python-pptx exposes group children in the group's internal coordinate
    system. This helper composes each group transform so generation can insert
    top-level replacement pictures at the same visual position.
    """
    if path is None:
        return None
    shapes = getattr(slide_or_group, "shapes", slide_or_group)
    sx = sy = 1.0
    tx = ty = 0.0
    shape = None

    for depth, index in enumerate(path):
        try:
            shape = shapes[index]
        except (IndexError, TypeError, KeyError):
            return None

        left, top, width, height = _box_tuple(shape)
        abs_box = (
            int(round(tx + left * sx)),
            int(round(ty + top * sy)),
            int(round(width * sx)),
            int(round(height * sy)),
        )
        if depth == len(path) - 1:
            return abs_box
        transform = _group_child_transform(shape)
        if transform is None:
            return None
        child_sx, child_sy, child_off_x, child_off_y = transform
        tx = tx + left * sx - child_off_x * child_sx * sx
        ty = ty + top * sy - child_off_y * child_sy * sy
        sx *= child_sx
        sy *= child_sy
        shapes = shape.shapes
    return None


def shape_type_name(shape) -> str:
    try:
        st = shape.shape_type
        if st is None:
            return "UNKNOWN"
        return st.name if hasattr(st, "name") else str(st)
    except Exception:
        return "UNKNOWN"


def placeholder_info(shape):
    try:
        if shape.is_placeholder:
            return shape.placeholder_format.type.name if shape.placeholder_format.type else "PLACEHOLDER"
    except Exception:
        pass
    return None


def picture_fill_info(shape):
    """Detect bitmap fills on non-picture shapes.

    Templates commonly crop photos into circles or other shapes. python-pptx
    exposes those as AUTO_SHAPE with `a:blipFill`, not as PICTURE, but they are
    semantically image slots.
    """
    if qn is None:
        return None
    try:
        blip = shape._element.find(".//" + qn("a:blip"))  # noqa: SLF001
        if blip is None:
            return None
        return {"type": "blipFill", "embed": blip.get(qn("r:embed"))}
    except Exception:
        return None


def geometry_from_box(left: int, top: int, width: int, height: int,
                      slide_w: int, slide_h: int) -> dict:
    slide_area = max(1, slide_w * slide_h)
    cx = left + width / 2
    cy = top + height / 2
    on_canvas = not (
        left + width <= 0 or top + height <= 0 or left >= slide_w or top >= slide_h
    )
    return {
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "left_in": round(emu_to_inch(left), 3),
        "top_in": round(emu_to_inch(top), 3),
        "width_in": round(emu_to_inch(width), 3),
        "height_in": round(emu_to_inch(height), 3),
        "area_ratio": round((width * height) / slide_area, 4),
        "cx_ratio": round(cx / slide_w, 4) if slide_w else 0.0,
        "cy_ratio": round(cy / slide_h, 4) if slide_h else 0.0,
        "aspect": round(width / height, 3) if height else 0.0,
        "on_canvas": on_canvas,
    }


def geometry_of(shape, slide_w: int, slide_h: int) -> dict:
    return geometry_from_box(*_box_tuple(shape), slide_w, slide_h)


def absolute_geometry_of_path(slide_or_group, path, slide_w: int, slide_h: int) -> dict | None:
    box = absolute_box_for_path(slide_or_group, path)
    if box is None:
        return None
    return geometry_from_box(*box, slide_w, slide_h)


def _enum_name(value):
    return value.name if hasattr(value, "name") else str(value)


def _color_format_info(color):
    try:
        if color is None or color.type is None:
            return None
        out = {"type": _enum_name(color.type)}
        try:
            if color.rgb is not None:
                out["rgb"] = str(color.rgb)
        except Exception:
            pass
        try:
            if color.theme_color is not None:
                out["theme_color"] = _enum_name(color.theme_color)
        except Exception:
            pass
        try:
            out["brightness"] = color.brightness
        except Exception:
            pass
        return out
    except Exception:
        pass
    return None


def _run_color(run):
    try:
        return _color_format_info(run.font.color)
    except Exception:
        pass
    return None


def _enum_or_none(value):
    if value is None:
        return None
    return value.name if hasattr(value, "name") else str(value)


def _run_style_info(run) -> dict:
    """Return concrete run style for mixed-format text boxes."""
    out = {"text": run.text}
    try:
        if run.font.name:
            out["font_name"] = run.font.name
    except Exception:
        pass
    try:
        if run.font.size is not None:
            out["font_pt"] = run.font.size.pt
    except Exception:
        pass
    try:
        if run.font.bold is not None:
            out["bold"] = bool(run.font.bold)
    except Exception:
        pass
    try:
        if run.font.italic is not None:
            out["italic"] = bool(run.font.italic)
    except Exception:
        pass
    color = _run_color(run)
    if color:
        out["color"] = color
    return out


def _paragraph_style_info(para) -> dict:
    return {
        "text": para.text,
        "align": _enum_or_none(getattr(para, "alignment", None)),
        "runs": [_run_style_info(run) for run in para.runs],
    }


def _fill_info(fill):
    try:
        if fill is None or fill.type is None:
            return None
        out = {"type": _enum_name(fill.type)}
        try:
            out["transparency"] = fill.transparency
        except Exception:
            pass
        try:
            color = _color_format_info(fill.fore_color)
            if color:
                out["color"] = color
        except Exception:
            pass
        return out
    except Exception:
        pass
    return None


def _line_info(shape):
    try:
        line = shape.line
    except Exception:
        return None
    out = {}
    try:
        color = _color_format_info(line.color)
        if color:
            out["color"] = color
    except Exception:
        pass
    try:
        if line.width is not None:
            out["width"] = int(line.width)
    except Exception:
        pass
    try:
        out["transparency"] = line.transparency
    except Exception:
        pass
    return out or None


def visual_info(shape) -> dict | None:
    """Best-effort non-text visual style for fills, lines, and icons."""
    out = {}
    fill = _fill_info(getattr(shape, "fill", None))
    if fill:
        out["fill"] = fill
    line = _line_info(shape)
    if line:
        out["line"] = line
    return out or None


def text_info(shape) -> dict | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    tf = shape.text_frame
    full = tf.text.strip()
    if not full and len(tf.paragraphs) <= 1:
        # Empty text frame: still report (some templates use empty title boxes).
        pass
    sizes: list[float] = []
    bolds: list[bool] = []
    colors: list[dict] = []
    fonts: list[str] = []
    aligns: list = []
    paragraphs: list[dict] = []
    first_run_style: dict | None = None
    for para in tf.paragraphs:
        if para.alignment is not None:
            aligns.append(para.alignment)
        paragraphs.append(_paragraph_style_info(para))
        for run in para.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
            if run.font.name:
                fonts.append(run.font.name)
            if run.font.bold:
                bolds.append(True)
            col = _run_color(run)
            if col:
                colors.append(col)
            if first_run_style is None and run.text.strip():
                first_run_style = _run_style_info(run)
    if first_run_style is None:
        for para in paragraphs:
            for run in para.get("runs", []):
                if run:
                    first_run_style = run
                    break
            if first_run_style is not None:
                break
    align = aligns[0].name if aligns and hasattr(aligns[0], "name") else (str(aligns[0]) if aligns else None)
    font_name = (first_run_style or {}).get("font_name")
    if font_name is None and fonts:
        font_name = Counter(fonts).most_common(1)[0][0]
    return {
        "sample": full.replace("\n", " | ")[:60],
        "full": full,
        "char_count": len(full.replace("\n", "")),
        "n_paragraphs": len([p for p in tf.paragraphs if p.text.strip()]),
        "max_font_pt": max(sizes) if sizes else None,
        "min_font_pt": min(sizes) if sizes else None,
        "any_bold": bool(bolds),
        "font_name": font_name,
        "bold": (first_run_style or {}).get("bold"),
        "italic": (first_run_style or {}).get("italic"),
        "color": colors[0] if colors else None,
        "align": align,
        "paragraphs": paragraphs,
        "vertical": _is_vertical(tf),
    }


def _is_vertical(text_frame) -> bool:
    """Best-effort: read the bodyPr 'vert' attribute (vertical CJK titles)."""
    try:
        body = text_frame._txBody  # noqa: SLF001
        bodyPr = body.find("{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr")
        if bodyPr is not None:
            vert = bodyPr.get("vert")
            return bool(vert) and vert not in ("horz",)
    except Exception:
        pass
    return False
