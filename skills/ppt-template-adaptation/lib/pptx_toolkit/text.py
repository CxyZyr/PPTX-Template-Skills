"""Text writing helpers.

Extracted from scripts/generate_strategy_ppt.py (set_text, set_mixed_text,
set_run_style) and scripts/generate_strategy_ppt_ai_tech.py (replace_first_run_text).

Generalized: nothing here hardcodes a brand font or color. When a style argument
is None we inherit it from the template's existing first run, so a shape keeps
its template look unless the caller overrides it.
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


def _is_cjk_text(text: str) -> bool:
    """Check if text contains CJK (Chinese/Japanese/Korean) characters."""
    if not text:
        return False
    for char in text:
        code = ord(char)
        # CJK Unified Ideographs: 4E00-9FFF
        # CJK Extension A: 3400-4DBF
        # Hiragana: 3040-309F, Katakana: 30A0-30FF
        if (0x4E00 <= code <= 0x9FFF or
            0x3400 <= code <= 0x4DBF or
            0x3040 <= code <= 0x30FF):
            return True
    return False


def _is_vertical_text_frame(text_frame) -> bool:
    """Check if text_frame is set to vertical orientation."""
    try:
        # Check via XML element (more reliable)
        if hasattr(text_frame, '_element'):
            elem = text_frame._element
            bodyPr = elem.find('.//{*}bodyPr')
            if bodyPr is not None:
                vert = bodyPr.get('vert')
                # Various vertical text modes:
                # - 'vert', 'vert270': standard vertical
                # - 'wordArtVert': WordArt vertical
                # - 'eaVert': East Asian vertical (commonly used for Chinese)
                # - 'mongolianVert': Mongolian vertical
                return vert in ('vert', 'vert270', 'wordArtVert', 'eaVert', 'mongolianVert')
    except Exception:
        pass
    return False


def _set_horizontal_text_frame(text_frame) -> None:
    """Set text_frame to horizontal orientation."""
    try:
        if hasattr(text_frame, '_element'):
            elem = text_frame._element
            bodyPr = elem.find('.//{*}bodyPr')
            if bodyPr is not None:
                # Remove vert attribute to make it horizontal
                if 'vert' in bodyPr.attrib:
                    del bodyPr.attrib['vert']
    except Exception:
        pass


def _theme_color(value):
    if value is None:
        return None
    if isinstance(value, MSO_THEME_COLOR):
        return value
    name = str(value)
    try:
        return getattr(MSO_THEME_COLOR, name)
    except AttributeError:
        return None


def _apply_color(run, color) -> None:
    if color is None:
        return
    if isinstance(color, dict):
        rgb = color.get("rgb")
        if rgb:
            run.font.color.rgb = RGBColor.from_string(str(rgb).lstrip("#"))
            return
        theme = _theme_color(color.get("theme_color"))
        if theme is not None:
            run.font.color.theme_color = theme
            try:
                if color.get("brightness") is not None:
                    run.font.color.brightness = color["brightness"]
            except Exception:
                pass
            return
    if isinstance(color, RGBColor):
        run.font.color.rgb = color
    else:
        run.font.color.rgb = RGBColor.from_string(str(color).lstrip("#"))


def _read_run_color(run):
    try:
        color = run.font.color
        if color is None or color.type is None:
            return None
        try:
            if color.rgb is not None:
                return {"rgb": str(color.rgb)}
        except Exception:
            pass
        try:
            if color.theme_color is not None:
                out = {"theme_color": color.theme_color.name}
                try:
                    out["brightness"] = color.brightness
                except Exception:
                    pass
                return out
        except Exception:
            pass
    except Exception:
        pass
    return None


def set_run_style(run, font_name=None, size=None, *, bold=None, color=None) -> None:
    if font_name is not None:
        run.font.name = font_name
    if size is not None:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    _apply_color(run, color)


def _template_run(shape):
    tf = shape.text_frame
    para = tf.paragraphs[0] if tf.paragraphs else None
    return para.runs[0] if (para and para.runs) else None


def set_text(
    shape,
    lines,
    *,
    font_name=None,
    size=None,
    bold=None,
    color=None,
    align=None,
    valign=None,
    line_spacing=1.1,
    space_after=2,
    margins=None,
) -> None:
    """Replace a shape's text with one paragraph per line.

    Unspecified style (font_name/size/color/align) is inherited from the
    template's first run so the shape keeps its original styling by default.
    """
    if isinstance(lines, str):
        lines = [lines]

    text_frame = shape.text_frame
    template_run = _template_run(shape)
    template_para = text_frame.paragraphs[0] if text_frame.paragraphs else None

    if font_name is None and template_run is not None:
        font_name = template_run.font.name
    if size is None and template_run is not None:
        size = template_run.font.size
    if color is None and template_run is not None:
        color = _read_run_color(template_run)
    if align is None and template_para is not None:
        align = template_para.alignment

    text_frame.clear()
    text_frame.word_wrap = True
    if valign is not None:
        text_frame.vertical_anchor = valign
    if margins:
        left, top, right, bottom = margins
        text_frame.margin_left = Pt(left)
        text_frame.margin_top = Pt(top)
        text_frame.margin_right = Pt(right)
        text_frame.margin_bottom = Pt(bottom)

    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        if align is not None:
            paragraph.alignment = align
        paragraph.space_after = Pt(space_after)
        paragraph.line_spacing = line_spacing
        run = paragraph.add_run()
        run.text = line
        set_run_style(run, font_name, size, bold=bold, color=color)


def set_mixed_text(shape, segments, *, align=PP_ALIGN.LEFT, valign=None, margins=None,
                   line_spacing=1.12) -> None:
    """Write multiple lines with per-line styling, inside a single paragraph
    (line breaks instead of new paragraphs — keeps tight spacing).

    `segments` is a list of dicts: {text, font_name?, size?, bold?, color?}.
    """
    text_frame = shape.text_frame
    template_run = _template_run(shape)
    fallback_font = template_run.font.name if template_run else None
    fallback_size = template_run.font.size if template_run else None

    text_frame.clear()
    text_frame.word_wrap = True
    if valign is not None:
        text_frame.vertical_anchor = valign
    if margins:
        left, top, right, bottom = margins
        text_frame.margin_left = Pt(left)
        text_frame.margin_top = Pt(top)
        text_frame.margin_right = Pt(right)
        text_frame.margin_bottom = Pt(bottom)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    paragraph.space_after = Pt(0)
    for index, segment in enumerate(segments):
        if index > 0:
            paragraph.add_line_break()
        run = paragraph.add_run()
        run.text = segment["text"]
        set_run_style(
            run,
            segment.get("font_name", fallback_font),
            segment.get("size", fallback_size),
            bold=segment.get("bold"),
            color=segment.get("color"),
        )


def replace_first_run_text(shape, text: str) -> None:
    """Overwrite only the first run's text (preserves the run's exact styling).

    Useful for strong template-dependent elements (cover title, ring center
    value, badge number) where rebuilding the text frame would lose styling.
    """
    text_frame = shape.text_frame
    if not text_frame.paragraphs:
        return
    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = text
