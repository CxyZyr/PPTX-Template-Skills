"""apply_content_plan.py — fill a PPTX template from a content_plan.json against a
parsed spec.json, reusing lib/pptx_toolkit. The generic generation entry point.

Pipeline:  template.pptx + spec.json + content_plan.json  ->  filled deck.pptx (+renders)

Usage:
  # full build (fills every page, drops slides not in keep_slides, optional render):
  python3 apply_content_plan.py --plan workspace/plans/gcl_strategy.content_plan.json --render

  # incremental, page-by-page work (fills only these slides, keeps the full deck so
  # original indices stay valid while you iterate; does NOT drop unused slides):
  python3 apply_content_plan.py --plan ...content_plan.json --pages 3,4 --render

Paths inside the plan (spec / output / asset_dir / logo / images) are resolved relative
to the current working directory (run from the repo root).
"""
from __future__ import annotations

import argparse
import json
import os
import pathlib
import sys

SKILL_ROOT = pathlib.Path(__file__).resolve().parents[1]
REPO_ROOT = pathlib.Path(__file__).resolve().parents[3]
sys.path.insert(0, str(REPO_ROOT))
sys.path.insert(0, str(SKILL_ROOT))

from pptx import Presentation  # noqa: E402

from lib.pptx_toolkit import apply_page, deck, render  # noqa: E402
from tavily_image_service import adapt_plan_images, load_dotenv  # noqa: E402
from validate_content_plan import print_report, validate_plan  # noqa: E402


def default_brand(spec: dict) -> dict:
    """Derive brand defaults from the template theme (overridden by plan['brand'])."""
    theme = spec["template"].get("theme", {})
    fonts = theme.get("common_fonts") or []
    colors = theme.get("common_colors") or []
    brand: dict = {}
    if fonts:
        brand["title_font"] = fonts[0]
    if colors:
        hexv = colors[0].lstrip("#")
        brand["accent"] = hexv
    return brand


def resolve_asset(path, asset_dir: pathlib.Path):
    """Resolve a file asset: as-is if it exists, else under asset_dir."""
    if not path:
        return path
    p = pathlib.Path(path)
    if p.is_absolute() or p.exists():
        return str(p)
    candidate = asset_dir / path
    return str(candidate) if candidate.exists() else str(p)


def resolve_image_item(item, asset_dir: pathlib.Path):
    """Resolve image path fields while preserving slot order and fallback items."""
    if isinstance(item, dict):
        out = dict(item)
        if out.get("path"):
            out["path"] = resolve_asset(out["path"], asset_dir)
        return out
    if isinstance(item, str):
        return resolve_asset(item, asset_dir) if item else ""
    return item


def has_unresolved_tavily_images(plan: dict, selected_pages: set[int] | None) -> bool:
    for page in plan.get("pages", []):
        if selected_pages is not None and page.get("slide") not in selected_pages:
            continue
        for item in page.get("images") or []:
            if not isinstance(item, dict):
                continue
            provider = item.get("provider") or item.get("source")
            if provider in ("tavily", "web") and not item.get("path") and item.get("fallback") != "never":
                return True
    return False


def resolve_web_images(plan: dict, asset_dir: pathlib.Path,
                       selected_pages: set[int] | None) -> dict:
    if not has_unresolved_tavily_images(plan, selected_pages):
        return plan
    api_key = os.environ.get("TAVILY_API_KEY", "")
    if not api_key:
        print("image provider: Tavily requested but TAVILY_API_KEY is not set; keeping original images")
        return plan
    print("image provider: resolving Tavily image slots")
    try:
        return adapt_plan_images(
            plan,
            api_key=api_key,
            asset_dir=asset_dir,
            selected_pages=selected_pages,
            query_prefix="",
            max_results=8,
            search_depth="basic",
            min_width=640,
            min_height=360,
            aspect_tolerance=2.2,
            max_bytes=8_000_000,
            force=False,
            dry_run=False,
        )
    except Exception as exc:
        print(f"image provider: Tavily failed ({exc}); keeping original images")
        return plan


def print_image_provider_report(plan: dict, selected_pages: set[int] | None) -> None:
    rows = []
    for page in plan.get("pages", []):
        if selected_pages is not None and page.get("slide") not in selected_pages:
            continue
        for i, item in enumerate(page.get("images") or [], start=1):
            if not isinstance(item, dict):
                continue
            provider = item.get("provider") or item.get("source")
            if provider not in ("tavily", "web"):
                continue
            status = item.get("status") or ("resolved" if item.get("path") else "keep_original")
            path = item.get("path") or ""
            detail = f"slide {page.get('slide'):02d} image {i}: {status}"
            if path:
                detail += f" -> {path}"
            if item.get("page_url"):
                detail += f" source={item.get('page_url')}"
            elif item.get("url"):
                detail += f" source={item.get('url')}"
            elif item.get("fallback"):
                detail += f" ({item.get('fallback')})"
            rows.append(detail)
    if rows:
        print("image provider report:")
        for row in rows:
            print(f"  {row}")


def main() -> None:
    load_dotenv()
    ap = argparse.ArgumentParser(description="Fill a template from a content_plan + spec.")
    ap.add_argument("--plan", required=True, help="content_plan.json path")
    ap.add_argument("--pages", default=None,
                    help="comma-separated original slide indices to fill (incremental mode)")
    ap.add_argument("--render", action="store_true", help="export PDF + per-slide PNGs after saving")
    ap.add_argument("--no-clear-unfilled", action="store_true",
                    help="keep template placeholder text in slots you didn't fill")
    ap.add_argument("--skip-global-text-check", action="store_true",
                    help="skip contents/section global text coherence validation")
    ap.add_argument("--skip-image-resolve", action="store_true",
                    help="skip web image providers and keep unresolved image slots original")
    args = ap.parse_args()

    plan = json.loads(pathlib.Path(args.plan).read_text(encoding="utf-8"))
    spec = json.loads(pathlib.Path(plan["spec"]).read_text(encoding="utf-8"))

    asset_dir = pathlib.Path(plan.get("asset_dir", "."))
    output_path = pathlib.Path(plan["output"])
    output_path.parent.mkdir(parents=True, exist_ok=True)
    brand = {**default_brand(spec), **(plan.get("brand") or {})}
    slide_width = spec["template"]["slide_width_emu"]
    slide_height = spec["template"]["slide_height_emu"]

    fill_plan_by_index = {s["index"]: s["fill_plan"] for s in spec["slides"]}

    partial = None
    if args.pages:
        partial = {int(x) for x in args.pages.split(",") if x.strip() != ""}

    if not args.skip_global_text_check:
        issues = validate_plan(plan, selected_pages=partial, spec=spec)
        print_report(issues)
        if any(issue.level == "error" for issue in issues):
            raise SystemExit("content_plan global text coherence check failed")

    if not args.skip_image_resolve:
        plan = resolve_web_images(plan, asset_dir, partial)
        print_image_provider_report(plan, partial)

    prs = Presentation(spec["template"]["path"])
    slides = list(prs.slides)

    print(f"template : {spec['template']['path']}")
    print(f"output   : {output_path}")
    print(f"brand    : {brand}")
    print(f"mode     : {'incremental ' + str(sorted(partial)) if partial else 'full build'}")
    print("-" * 72)

    total_warnings = 0
    for page in plan["pages"]:
        index = page["slide"]
        if partial is not None and index not in partial:
            continue
        if index not in fill_plan_by_index:
            print(f"slide {index:02d}  [skip] not in spec")
            continue
        # resolve file assets (icons stay as bootstrap names, handled by apply_page)
        page = dict(page)
        if page.get("logo"):
            page["logo"] = resolve_asset(page["logo"], asset_dir)
        if "images" in page:
            page["images"] = [resolve_image_item(p, asset_dir) for p in page["images"]]

        report = apply_page(
            slides[index], fill_plan_by_index[index], page,
            asset_dir=asset_dir, brand=brand,
            slide_width=slide_width, slide_height=slide_height,
            clear_unfilled=not args.no_clear_unfilled,
        )
        filled = ", ".join(f"{k}={v}" for k, v in report["filled"].items()) or "(nothing)"
        print(f"slide {index:02d}  [{report.get('role',''):15s}] {filled}")
        for w in report["warnings"]:
            total_warnings += 1
            print(f"           ! {w}")

    # Drop unused slides only on a full build (incremental mode keeps indices stable).
    if partial is None and "keep_slides" in plan:
        keep = set(plan["keep_slides"])
        drop = [i for i in range(len(slides)) if i not in keep]
        if drop:
            deck.delete_slides(prs, drop)
            print("-" * 72)
            print(f"dropped slides: {sorted(drop)}  ->  {len(prs.slides)} kept")

    prs.save(str(output_path))
    print("-" * 72)
    print(f"saved -> {output_path}   (warnings: {total_warnings})")

    if args.render:
        renders = render.render_deck(output_path, output_path.parent / "renders")
        print(f"rendered {len(renders)} page(s) -> {output_path.parent / 'renders'}")


if __name__ == "__main__":
    main()
